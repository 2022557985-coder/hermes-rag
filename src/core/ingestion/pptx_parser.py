"""PPTX parser using python-pptx."""

from pathlib import Path
from typing import Dict, Any

from .parser_factory import BaseParser


class PptxParser(BaseParser):
    """Parse PPTX presentations."""

    def parse(self, source: str) -> Dict[str, Any]:
        """Parse a PPTX file.

        Args:
            source: Path to PPTX file.

        Returns:
            dict with text, tables, metadata.
        """
        from pptx import Presentation

        prs = Presentation(source)
        slides_text = []
        tables = []
        metadata = {
            "source": Path(source).name,
            "type": "pptx",
            "total_slides": len(prs.slides),
        }

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title
                            slide_title = text
                            slide_content.append(f"## {text}")
                        else:
                            slide_content.append(text)

                if shape.has_table:
                    md_table = self._table_to_markdown(shape.table)
                    if md_table:
                        tables.append({
                            "slide": slide_num,
                            "content": md_table,
                        })
                        slide_content.append(md_table)

            slides_text.append({
                "slide": slide_num,
                "title": slide_title,
                "content": "\n".join(slide_content),
            })

        full_text = "\n\n".join(
            f"Slide {s['slide']}: {s['title']}\n{s['content']}"
            for s in slides_text
        )

        return {
            "text": full_text,
            "tables": tables,
            "metadata": metadata,
        }